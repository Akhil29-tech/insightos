from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

TEAL = RGBColor(0x1D, 0x9E, 0x75)
DARK = RGBColor(0x1E, 0x1E, 0x1C)
GRAY = RGBColor(0x6E, 0x6E, 0x69)

def build_presentation(entries: list) -> bytes:
    """entries: list of {question, narration, sql, results}"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1.2))
    tf = title_box.text_frame
    tf.text = "InsightOS — Retail Analysis"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(0.6))
    sf = subtitle_box.text_frame
    sf.text = f"{len(entries)} findings from this session"
    sf.paragraphs[0].font.size = Pt(18)
    sf.paragraphs[0].font.color.rgb = GRAY

    # One slide per finding
    for entry in entries:
        slide = prs.slides.add_slide(blank_layout)

        q_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.7))
        qf = q_box.text_frame
        qf.word_wrap = True
        qf.text = entry.get("question", "")
        qf.paragraphs[0].font.size = Pt(22)
        qf.paragraphs[0].font.bold = True
        qf.paragraphs[0].font.color.rgb = DARK

        n_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12), Inches(0.8))
        nf = n_box.text_frame
        nf.word_wrap = True
        nf.text = entry.get("narration", "")
        nf.paragraphs[0].font.size = Pt(16)
        nf.paragraphs[0].font.color.rgb = TEAL

        results = entry.get("results", [])[:8]
        if results:
            cols = list(results[0].keys())
            rows = len(results) + 1
            table_shape = slide.shapes.add_table(
                rows, len(cols), Inches(0.6), Inches(2.2), Inches(12), Inches(0.4 * rows)
            )
            table = table_shape.table
            for c, col in enumerate(cols):
                cell = table.cell(0, c)
                cell.text = str(col)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)
            for r, row in enumerate(results, start=1):
                for c, col in enumerate(cols):
                    cell = table.cell(r, c)
                    cell.text = str(row.get(col, ""))
                    cell.text_frame.paragraphs[0].font.size = Pt(11)

        sql_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.65), Inches(12.2), Inches(0.6))
        sqf = sql_box.text_frame
        sqf.word_wrap = True
        raw_sql = " ".join(entry.get("sql", "").split())
        sqf.text = raw_sql[:180] + ("..." if len(raw_sql) > 180 else "")
        for para in sqf.paragraphs:
            para.font.size = Pt(9)
            para.font.color.rgb = GRAY

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
